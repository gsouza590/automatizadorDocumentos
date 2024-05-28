import os
import PySimpleGUI as sg
from pptx import Presentation
from utils import load_excel_data

def replace_text_in_honorarios(combo_value, pptx_file, output_path, input_path, objeto_proposta=None,
                                valor_de_honorarios=None, obs_honorario=None, parcelamento=None,
                               validade_da_proposta=None, prazo_de_execucao=None):
    excel_file_path = os.path.join(input_path, "NLcadastro.xlsm")
    ws, names = load_excel_data(excel_file_path)
    if not ws:
        sg.popup_error("Erro ao carregar o arquivo Excel.")
        return

    info = None
    for row in ws.iter_rows(min_row=2, max_col=ws.max_column, values_only=True):
        if row[5] == combo_value:
            info = {
                '#nome': row[5],
                'objetoproposta': objeto_proposta,
                'valordehonorarios': valor_de_honorarios,
                'obshonorario': obs_honorario,
                'parchonorarios': parcelamento,
                'validadeproposta': validade_da_proposta,
                'prazoexecucao': prazo_de_execucao
            }
            # Remove any items where the value is None
            info = {k: v for k, v in info.items() if v is not None}
            break
    if not info:
        sg.popup("Nome não encontrado.")
        return

    prs = Presentation(pptx_file)
    text_replaced = False

    def replace_in_shape(shape, info):
        nonlocal text_replaced
        if hasattr(shape, "text_frame") and shape.text_frame is not None:
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    for placeholder, value in info.items():
                        if placeholder in run.text:
                            run.text = run.text.replace(placeholder, str(value))
                            text_replaced = True
            for placeholder, value in info.items():
                if placeholder in shape.text_frame.text:
                    shape.text_frame.text = shape.text_frame.text.replace(placeholder, str(value))
                    text_replaced = True
        elif hasattr(shape, "table") and shape.table is not None:
            for row in shape.table.rows:
                for cell in row.cells:
                    replace_in_shape(cell, info)
        elif hasattr(shape, "chart") and shape.chart is not None:
            chart = shape.chart
            for series in chart.series:
                for point in series.points:
                    for placeholder, value in info.items():
                        if placeholder in point.data_label.text_frame.text:
                            point.data_label.text_frame.text = point.data_label.text_frame.text.replace(placeholder, str(value))
                            text_replaced = True

    for slide in prs.slides:
        for shape in slide.shapes:
            replace_in_shape(shape, info)

    if text_replaced:
        output_file = os.path.join(output_path, f"Honorarios_{combo_value.replace(' ', '_')}.pptx")
        prs.save(output_file)
        sg.popup("Criado com sucesso!")
    else:
        sg.popup_error("Nenhuma substituição de texto foi realizada.")
