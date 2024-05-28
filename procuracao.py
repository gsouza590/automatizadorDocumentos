import os
import PySimpleGUI as sg
from pptx import Presentation
from utils import load_excel_data

def replace_text_in_procuracao(combo_value, finalidade, pptx_file, output_path, input_path):
    excel_file_path = os.path.join(input_path, "NLcadastro.xlsm")
    ws, names = load_excel_data(excel_file_path)
    if not ws:
        return

    info = None
    for row in ws.iter_rows(min_row=2, max_col=ws.max_column, values_only=True):
        if row[5] == combo_value:
            info = {
                '#nome': row[5],
                '#cpf': row[9],
                '#endereco': f"{row[15]} {row[16]}, {row[17]} - CEP: {row[18]}, {row[19]} - {row[20]}",
                '#email': row[1],
                '#finalidade': finalidade
            }
            break
    if not info:
        sg.popup("Nome não encontrado.")
        return

    prs = Presentation(pptx_file)
    text_replaced = False
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        for placeholder, value in info.items():
                            if placeholder in run.text:
                                run.text = run.text.replace(placeholder, str(value))
                                text_replaced = True
    if text_replaced:
        output_file = os.path.join(output_path, f"Procuracao_{combo_value.replace(' ', '_')}.pptx")
        prs.save(output_file)
        sg.popup("Apresentação gerada com sucesso!")
    else:
        sg.popup_error("Nenhuma substituição de texto foi realizada.")
